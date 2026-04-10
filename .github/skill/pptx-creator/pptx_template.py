"""
PPTX Template — Microsoft-Branded PowerPoint Generator
Uses python-pptx to create professional 10-slide presentations.

Usage:
    python3 pptx_template.py

Requirements:
    pip install python-pptx

Customization:
    Replace the {{variables}} and sample content with real data.
    The agent should modify this template for each presentation request.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from datetime import date
import os

# === CONFIG ===
TITLE = "{{title}}"
SUBTITLE = "{{subtitle}}"
AUTHOR = "{{author_name}}"
CLIENT = "{{client_name}}"
VERSION = "1.0.0"
DATE = date.today().isoformat()

# === MICROSOFT BRAND COLORS ===
BLUE = RGBColor(0x00, 0x78, 0xD4)
RED = RGBColor(0xF2, 0x50, 0x22)
GREEN = RGBColor(0x7F, 0xBA, 0x00)
LIGHT_BLUE = RGBColor(0x00, 0xA4, 0xEF)
YELLOW = RGBColor(0xFF, 0xB9, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_PRIMARY = RGBColor(0x32, 0x31, 0x30)
TEXT_SECONDARY = RGBColor(0x60, 0x5E, 0x5C)
ALT_ROW = RGBColor(0xF3, 0xF2, 0xF1)


def add_4color_bar(slide, top, width, height=Emu(36000)):
    """Add Microsoft 4-color bar to a slide."""
    colors = [RED, GREEN, LIGHT_BLUE, YELLOW]
    segment_width = width // 4
    for i, color in enumerate(colors):
        shape = slide.shapes.add_shape(
            1, Emu(i * segment_width), top, Emu(segment_width), height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()


def set_text(shape, text, font_size=18, bold=False, color=TEXT_PRIMARY, alignment=PP_ALIGN.LEFT):
    """Set text in a shape with formatting."""
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Segoe UI'
    return tf


def add_slide(prs, layout_index=5):
    """Add a blank slide."""
    layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(layout)


def add_text_box(slide, left, top, width, height, text, font_size=16, bold=False, color=TEXT_PRIMARY, alignment=PP_ALIGN.LEFT):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Segoe UI'
    return txBox


def add_branded_table(slide, left, top, width, height, headers, rows):
    """Add a Microsoft-branded table to a slide."""
    table_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), left, top, width, height
    )
    table = table_shape.table

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = 'Segoe UI'
                run.font.size = Pt(10)
                run.font.bold = True
                run.font.color.rgb = WHITE

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ALT_ROW
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = 'Segoe UI'
                    run.font.size = Pt(9)
                    run.font.color.rgb = TEXT_PRIMARY

    return table


def add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


def create_presentation():
    """Generate the complete 10-slide branded presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # === SLIDE 1: EXECUTIVE COVER ===
    slide = add_slide(prs)
    add_4color_bar(slide, Emu(0), slide_width, Emu(72000))
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                 TITLE, font_size=36, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4), Inches(11), Inches(0.8),
                 SUBTITLE, font_size=18, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                 f"{AUTHOR} | {DATE} | v{VERSION}", font_size=11, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
                 "Microsoft Confidential", font_size=9, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
    add_speaker_notes(slide, f"Cover slide for {TITLE}. Welcome the audience and set context.")

    # === SLIDE 2: AGENDA ===
    slide = add_slide(prs)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "Agenda", font_size=28, bold=True, color=BLUE)
    agenda_items = [
        "1. Context & Challenge",
        "2. Solution Overview",
        "3. Technical Deep-Dive",
        "4. Implementation & Demo",
        "5. Benefits & ROI",
        "6. Timeline & Next Steps",
    ]
    for i, item in enumerate(agenda_items):
        add_text_box(slide, Inches(1), Inches(1.5 + i * 0.7), Inches(10), Inches(0.6),
                     item, font_size=18, color=TEXT_PRIMARY)
    add_speaker_notes(slide, "Walk through the agenda. Set expectations for timing.")

    # === SLIDES 3-9: CONTENT SLIDES ===
    content_slides = [
        ("Context & Challenge", "Describe the current state, pain points, and why change is needed."),
        ("Solution Overview", "Present the solution architecture and approach at a high level."),
        ("Technical Deep-Dive 1", "Component breakdown, technologies, and integration details."),
        ("Technical Deep-Dive 2", "Workflows, processes, and data flows."),
        ("Demo / Visual", "Screenshots, architecture diagrams, or live demo walkthrough."),
        ("Benefits & Value", "ROI metrics, KPIs, and comparison tables."),
        ("Timeline & Roadmap", "Phased implementation plan with milestones."),
    ]
    for title, description in content_slides:
        slide = add_slide(prs)
        add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                     title, font_size=28, bold=True, color=BLUE)
        add_text_box(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(5),
                     description, font_size=16, color=TEXT_PRIMARY)
        add_speaker_notes(slide, f"Speaker notes for: {title}")

    # === SLIDE 10: NEXT STEPS ===
    slide = add_slide(prs)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "Next Steps", font_size=28, bold=True, color=BLUE)
    add_branded_table(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(3),
                      ['#', 'Action', 'Owner', 'Deadline'],
                      [
                          ['1', 'Action item 1', 'Owner', 'Date'],
                          ['2', 'Action item 2', 'Owner', 'Date'],
                          ['3', 'Action item 3', 'Owner', 'Date'],
                      ])
    add_text_box(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.5),
                 "Sources: [1] URL, [2] URL", font_size=9, color=TEXT_SECONDARY)
    add_speaker_notes(slide, "Close with clear next steps and call to action.")

    # === SAVE ===
    os.makedirs('output/pptx', exist_ok=True)
    filename = f"{TITLE.replace(' ', '_')}_v{VERSION}_{DATE}.pptx"
    filepath = f"output/pptx/{filename}"
    prs.save(filepath)
    print(f"Created: {filepath}")


if __name__ == '__main__':
    create_presentation()
